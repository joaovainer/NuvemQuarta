from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "apresentacao_mural_distribuido.pptx"

BG = RGBColor(246, 241, 232)
INK = RGBColor(29, 37, 40)
MUTED = RGBColor(88, 101, 105)
ACCENT = RGBColor(208, 95, 50)
GREEN = RGBColor(29, 122, 85)


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title, subtitle):
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.75))
    frame = title_box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.name = "Georgia"
    paragraph.font.size = Pt(31)
    paragraph.font.bold = True
    paragraph.font.color.rgb = INK

    sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(1.03), Inches(11.5), Inches(0.35))
    sub_frame = sub_box.text_frame
    sub_frame.clear()
    sub = sub_frame.paragraphs[0]
    sub.text = subtitle
    sub.font.name = "Verdana"
    sub.font.size = Pt(11)
    sub.font.color.rgb = MUTED


def add_card(slide, x, y, w, h, text, color=INK, fill=RGBColor(255, 252, 245)):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(220, 205, 181)
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = "Verdana"
    paragraph.font.size = Pt(14)
    paragraph.font.bold = True
    paragraph.font.color.rgb = color
    return shape


def add_bullets(slide, items, x, y, w, h, font_size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Verdana"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = INK
        paragraph.space_after = Pt(10)


def add_footer(slide):
    box = slide.shapes.add_textbox(Inches(0.58), Inches(7.05), Inches(12), Inches(0.25))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = "Mural de Ideias Distribuido | FastAPI + PostgreSQL + Redis + Worker | Deploy Render"
    paragraph.font.name = "Verdana"
    paragraph.font.size = Pt(8)
    paragraph.font.color.rgb = MUTED


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Arquitetura distribuida", "Premissas: servicos independentes, rede, fila, persistencia, cache e health checks.")

    cards = [
        (0.75, 2.05, "Navegador\nHTML/CSS/JS"),
        (3.15, 2.05, "API REST\nFastAPI stateless"),
        (5.55, 1.45, "PostgreSQL\npersistencia"),
        (5.55, 2.72, "Redis\nfila + cache"),
        (8.0, 2.05, "Worker\nprocessamento async"),
        (10.45, 2.05, "Usuario\nresultado na tela"),
    ]
    for x, y, text in cards:
        fill = RGBColor(255, 252, 245)
        color = INK
        if "Redis" in text:
            fill = RGBColor(255, 232, 220)
            color = ACCENT
        if "Worker" in text:
            fill = RGBColor(224, 242, 234)
            color = GREEN
        add_card(slide, x, y, 1.85, 0.82, text, color=color, fill=fill)

    add_bullets(
        slide,
        [
            "API grava a ideia no banco e envia o ID para a fila Redis.",
            "Worker consome a fila sem travar a resposta ao usuario.",
            "Cache reduz consultas repetidas e health checks mostram a operacao.",
        ],
        1.0,
        4.35,
        11.2,
        1.5,
        17,
    )
    add_footer(slide)


def slide_cloud(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Escolha da nuvem: Render", "Deploy publico usando Blueprint com Docker, web services, Postgres e Key Value/Redis.")

    add_card(slide, 0.9, 1.7, 3.65, 1.2, "1 arquivo\nrender.yaml", color=ACCENT)
    add_card(slide, 4.85, 1.7, 3.65, 1.2, "4 recursos\nAPI + Worker + DB + Redis", color=GREEN)
    add_card(slide, 8.8, 1.7, 3.65, 1.2, "URL publica\npara demonstracao", color=INK)

    add_bullets(
        slide,
        [
            "Render simplifica a publicacao: conecta GitHub e provisiona recursos pelo Blueprint.",
            "PostgreSQL e Key Value/Redis ficam gerenciados pela nuvem, sem instalar servidor manualmente.",
            "Servicos Docker deixam o mesmo projeto rodar localmente e em producao.",
            "Plano gratuito serve para demonstracao; em producao, basta trocar o plano e escalar.",
        ],
        1.05,
        3.6,
        11.2,
        2.2,
        17,
    )
    add_footer(slide)


def slide_benefits(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Beneficios da solucao", "Projeto pequeno, mas com os conceitos principais de sistemas distribuidos.")

    add_bullets(
        slide,
        [
            "Desacoplamento: a API responde rapido e o worker processa tarefas em segundo plano.",
            "Escalabilidade: API e worker podem ganhar mais instancias conforme a demanda.",
            "Confiabilidade: dados persistem no PostgreSQL e a fila Redis organiza o processamento.",
            "Observabilidade: health checks exibem banco, Redis, tamanho da fila e atividade do worker.",
            "Facil explicacao: fluxo simples de ponta a ponta, ideal para demonstrar ao professor.",
        ],
        0.95,
        1.75,
        11.5,
        3.4,
        18,
    )

    link_box = slide.shapes.add_textbox(Inches(0.95), Inches(5.7), Inches(11.4), Inches(0.6))
    frame = link_box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = "Link do sistema: substituir pela URL gerada pela Render apos o deploy"
    paragraph.font.name = "Verdana"
    paragraph.font.size = Pt(15)
    paragraph.font.bold = True
    paragraph.font.color.rgb = ACCENT
    add_footer(slide)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_architecture(prs)
    slide_cloud(prs)
    slide_benefits(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"PPT gerado em: {OUTPUT}")


if __name__ == "__main__":
    main()
